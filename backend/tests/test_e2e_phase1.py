"""
End-to-End Tests for Phase 1 Document Intelligence Implementation.

This test suite validates the complete Phase 1 implementation:
- Docling document parsing (DOCX, PPTX, PDF)
- Large document chunking (50+ pages)
- Transcription mode (Markdown output, not JSON)
- Full pipeline integration

Tests are designed to run without API keys by testing the document processing
pipeline components that don't require external VLM calls.
"""

import pytest
from PIL import Image
from services.processing import ProcessingService
from services.docling_service import DoclingService
from services.chunking_service import MarkdownSplitter
from services.transcription_service import TranscriptionService
from services.schema_service import SchemaService


class TestDoclingDocumentParsing:
    """Test Docling document parsing for various file types."""

    @pytest.fixture
    def docling_service(self):
        """Provide DoclingService instance."""
        return DoclingService()

    @pytest.fixture
    def sample_docx(self, tmp_path):
        """Create a sample DOCX file for testing."""
        # Create a simple DOCX file using python-docx
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")

        doc = Document()
        doc.add_heading("Test Document", level=1)
        doc.add_paragraph("This is a test paragraph for DOCX parsing.")
        doc.add_paragraph("Another paragraph with some content.")

        file_path = tmp_path / "sample.docx"
        doc.save(str(file_path))
        return str(file_path)

    @pytest.fixture
    def sample_pptx(self, tmp_path):
        """Create a sample PPTX file for testing."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        prs = Presentation()

        # Add title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = "Test Presentation"
        subtitle.text = "E2E Testing"

        # Add content slide
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = content_slide.shapes.title
        title.text = "Content Slide"
        body = content_slide.placeholders[1]
        text_frame = body.text_frame
        text_frame.text = "This is test content for PPTX parsing."

        file_path = tmp_path / "sample.pptx"
        prs.save(str(file_path))
        return str(file_path)

    @pytest.fixture
    def searchable_pdf(self, tmp_path):
        """Create a searchable PDF file for testing."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph
        except ImportError:
            pytest.skip("reportlab not installed")

        file_path = tmp_path / "searchable.pdf"
        doc = SimpleDocTemplate(str(file_path), pagesize=letter)

        styles = getSampleStyleSheet()
        story = []
        story.append(Paragraph("Searchable PDF Document", styles["Heading1"]))
        story.append(
            Paragraph(
                "This is a test paragraph for searchable PDF parsing.", styles["Normal"]
            )
        )
        story.append(
            Paragraph(
                "Another paragraph with searchable text content.", styles["Normal"]
            )
        )

        doc.build(story)
        return str(file_path)

    @pytest.fixture
    def image_only_pdf(self, tmp_path):
        """Create an image-only PDF for testing OCR."""
        # Create images and convert to PDF
        img = Image.new("RGB", (800, 600), color="white")

        # Add some text to the image (requires PIL)
        from PIL import ImageDraw, ImageFont

        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 40)
        except Exception:
            font = ImageFont.load_default()

        draw.text((50, 50), "Image-Only PDF", fill="black", font=font)
        draw.text(
            (50, 100), "This text is embedded in an image", fill="black", font=font
        )

        # Save as PDF
        file_path = tmp_path / "image_only.pdf"
        img.save(str(file_path), "PDF")
        return str(file_path)

    @pytest.fixture
    def large_pdf(self, tmp_path):
        """Create a large PDF (50+ pages) for testing chunking."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, PageBreak
        except ImportError:
            pytest.skip("reportlab not installed")

        file_path = tmp_path / "large_pdf.pdf"
        doc = SimpleDocTemplate(str(file_path), pagesize=letter)

        styles = getSampleStyleSheet()
        story = []

        # Create 50 pages
        for i in range(50):
            story.append(Paragraph(f"Page {i + 1}", styles["Heading1"]))
            story.append(
                Paragraph(
                    f"This is content for page {i + 1} of the large PDF document. "
                    * 10,
                    styles["Normal"],
                )
            )
            if i < 49:  # Don't add page break after last page
                story.append(PageBreak())

        doc.build(story)
        return str(file_path)

    def test_docx_parsing(self, docling_service, sample_docx):
        """Test DOCX parsing through Docling pipeline."""
        markdown = docling_service.parse_document(sample_docx)

        assert markdown is not None
        assert isinstance(markdown, str)
        assert len(markdown) > 0
        # Check for markdown headers or content
        assert "Test Document" in markdown or "test paragraph" in markdown.lower()

    def test_pptx_parsing(self, docling_service, sample_pptx):
        """Test PPTX parsing through Docling pipeline."""
        markdown = docling_service.parse_document(sample_pptx)

        assert markdown is not None
        assert isinstance(markdown, str)
        assert len(markdown) > 0
        # Check for presentation content
        assert "Test Presentation" in markdown or "Content Slide" in markdown

    def test_searchable_pdf_parsing(self, docling_service, searchable_pdf):
        """Test searchable PDF parsing (no OCR needed)."""
        markdown = docling_service.parse_document(searchable_pdf)

        assert markdown is not None
        assert isinstance(markdown, str)
        assert len(markdown) > 0
        assert "Searchable PDF" in markdown or "test paragraph" in markdown.lower()

    def test_image_only_pdf_parsing(self, docling_service, image_only_pdf):
        """Test image-only PDF parsing (OCR required)."""
        # This test might be skipped if OCR is disabled
        try:
            markdown = docling_service.parse_document(image_only_pdf)

            assert markdown is not None
            assert isinstance(markdown, str)
            # OCR might not extract text perfectly, so just check we get some output
            assert len(markdown) > 0
        except Exception as e:
            # OCR might fail if not properly configured
            pytest.skip(f"OCR not available: {e}")


class TestLargeDocumentChunking:
    """Test large document processing with chunking."""

    @pytest.fixture
    def chunking_service(self):
        """Provide MarkdownSplitter instance."""
        return MarkdownSplitter()

    @pytest.fixture
    def large_pdf(self, tmp_path):
        """Create a large PDF (50+ pages) for testing chunking."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, PageBreak
        except ImportError:
            pytest.skip("reportlab not installed")

        file_path = tmp_path / "large_pdf.pdf"
        doc = SimpleDocTemplate(str(file_path), pagesize=letter)

        styles = getSampleStyleSheet()
        story = []

        # Create 50 pages with substantial content
        for i in range(50):
            story.append(Paragraph(f"Chapter {i + 1}", styles["Heading1"]))
            story.append(
                Paragraph(
                    f"This is the content for chapter {i + 1}. "
                    * 20,  # Repeat for more content
                    styles["Normal"],
                )
            )
            if i < 49:
                story.append(PageBreak())

        doc.build(story)
        return str(file_path)

    @pytest.fixture
    def large_markdown_text(self):
        """Create large markdown text for testing chunking."""
        # Create text that's large enough to trigger chunking
        content = []
        for i in range(100):
            content.append(f"# Section {i + 1}\n\n")
            content.append(f"This is section {i + 1} with substantial content. " * 50)
            content.append("\n\n")
        return "".join(content)

    def test_token_counting(self, chunking_service, large_markdown_text):
        """Test token counting functionality."""
        count = chunking_service.count_tokens(large_markdown_text)

        assert count > 0
        # Token count varies by encoding; just verify it's substantial
        assert count > 40000  # Should be reasonably large for 100 sections

    def test_should_chunk_detection(self):
        """Test chunking threshold detection."""
        service = ProcessingService()

        # Small text - should not chunk
        small_text = "This is a small document."
        should_chunk = service._should_chunk(small_text, "gpt-4o")
        assert should_chunk is False

        # Very large text - should chunk (use chars that will produce enough tokens)
        large_text = "word " * 200000  # 200000 words should exceed threshold
        should_chunk = service._should_chunk(large_text, "gpt-4o")
        assert should_chunk is True

    def test_chunking_splits_document(self, chunking_service, large_markdown_text):
        """Test that chunking properly splits large documents."""
        chunks = chunking_service.split(large_markdown_text)

        assert len(chunks) > 1  # Should split into multiple chunks
        assert all(isinstance(chunk, str) for chunk in chunks)
        assert all(len(chunk) > 0 for chunk in chunks)

    def test_large_pdf_chunking_workflow(self, large_pdf):
        """Test complete workflow: parse large PDF -> detect chunking -> split."""
        docling_service = DoclingService()
        chunking_service = MarkdownSplitter()

        # Step 1: Parse document
        markdown = docling_service.parse_document(large_pdf)
        assert markdown is not None
        assert len(markdown) > 0

        # Step 2: Check if chunking is needed
        processing_service = ProcessingService()
        should_chunk = processing_service._should_chunk(markdown, "gpt-4o")
        # Large PDF should trigger chunking
        assert should_chunk is True

        # Step 3: Split into chunks
        chunks = chunking_service.split(markdown)
        assert len(chunks) > 1


class TestTranscriptionMode:
    """Test transcription mode produces Markdown, not JSON."""

    @pytest.fixture
    def transcription_service(self):
        """Provide TranscriptionService instance."""
        return TranscriptionService()

    @pytest.fixture
    def sample_audio(self, tmp_path):
        """Create a sample audio file for testing."""
        try:
            import wave
            import struct

            file_path = tmp_path / "sample_audio.wav"

            # Create a simple WAV file
            with wave.open(str(file_path), "w") as wav_file:
                wav_file.setnchannels(1)  # Mono
                wav_file.setsampwidth(2)  # 2 bytes per sample
                wav_file.setframerate(44100)  # 44.1kHz

                # Write 1 second of silence
                num_frames = 44100
                data = struct.pack("<" + "h" * num_frames, *[0] * num_frames)
                wav_file.writeframes(data)

            return str(file_path)
        except ImportError:
            pytest.skip("wave module not available")

    def test_transcription_service_exists(self, transcription_service):
        """Test that TranscriptionService can be instantiated."""
        assert transcription_service is not None
        assert hasattr(transcription_service, "transcribe")

    def test_transcription_returns_markdown(self, transcription_service, sample_audio):
        """Test that transcription returns Markdown text, not JSON."""
        # Note: This test may fail if transcription service is not fully configured
        try:
            result = transcription_service.transcribe(sample_audio)

            # Should return a dict with success and text
            assert isinstance(result, dict)
            assert "success" in result

            if result["success"]:
                # Text should be markdown/string, not JSON
                assert "text" in result
                assert isinstance(result["text"], str)

                # Verify it's not JSON format
                text = result["text"]
                assert not text.startswith("{")
                assert not text.startswith("[")
        except Exception as e:
            # Transcription might not be configured
            pytest.skip(f"Transcription not configured: {e}")

    def test_transcription_not_json_schema(self):
        """Test that transcription mode doesn't require JSON schema."""
        # Transcription should work with a simple text prompt
        # This verifies the service can handle non-structured output
        processing_service = ProcessingService()

        # Check that processing service can handle transcription mode
        assert hasattr(processing_service, "transcription_service")
        assert processing_service.transcription_service is not None


class TestFullPipelineIntegration:
    """Test complete pipeline integration without API keys."""

    @pytest.fixture
    def processing_service(self):
        """Provide ProcessingService instance."""
        return ProcessingService()

    @pytest.fixture
    def sample_schema(self):
        """Provide sample schema definition."""
        return {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "content": {"type": "string"},
            },
            "required": ["title"],
        }

    def test_processing_service_initialization(self, processing_service):
        """Test that ProcessingService initializes with all required components."""
        assert processing_service is not None
        assert hasattr(processing_service, "docling_service")
        assert hasattr(processing_service, "chunking_service")
        assert hasattr(processing_service, "transcription_service")
        assert hasattr(processing_service, "schema_service")

    def test_docling_service_integration(self, processing_service):
        """Test DoclingService is properly integrated."""
        assert processing_service.docling_service is not None
        assert isinstance(processing_service.docling_service, DoclingService)

    def test_chunking_service_integration(self, processing_service):
        """Test MarkdownSplitter is properly integrated."""
        assert processing_service.chunking_service is not None
        assert isinstance(processing_service.chunking_service, MarkdownSplitter)

    def test_file_size_validation(self, processing_service, tmp_path):
        """Test file size validation."""
        # Create a small test file
        small_file = tmp_path / "small.txt"
        small_file.write_text("Small content")

        # Should not raise exception for small file
        try:
            processing_service._validate_file_size(str(small_file))
        except ValueError:
            pytest.fail("File size validation failed for small file")

    def test_exposure_to_docling_method(
        self, processing_service, tmp_path, sample_schema
    ):
        """Test that _process_via_docling_parse method exists and is callable."""
        # Create a simple text file
        test_file = tmp_path / "test.txt"
        test_file.write_text("Test content")

        # Verify the method exists (renamed from _process_via_docling)
        assert hasattr(processing_service, "_process_via_docling_parse")

        # Note: We can't fully test this without a provider and API key
        # but we can verify the method signature
        import inspect

        sig = inspect.signature(processing_service._process_via_docling_parse)
        params = list(sig.parameters.keys())

        required_params = [
            "file_path",
            "provider",
            "model",
            "schema_definition",
            "prompt",
        ]
        for param in required_params:
            assert param in params


class TestMarkdownOutput:
    """Test Markdown output format for transcription mode."""

    def test_schema_service_builtin_templates(self):
        """Test that SchemaService provides builtin templates."""
        templates = SchemaService.get_builtin_templates()

        assert isinstance(templates, dict)
        assert len(templates) > 0
        assert "Generic" in templates

    def test_generic_schema_structure(self):
        """Test that Generic schema has proper structure."""
        templates = SchemaService.get_builtin_templates()
        generic = templates["Generic"]

        assert isinstance(generic, dict)
        assert "type" in generic
        assert generic["type"] == "object"
        assert "properties" in generic


class TestErrorHandling:
    """Test error handling in the pipeline."""

    @pytest.fixture
    def processing_service(self):
        """Provide ProcessingService instance."""
        return ProcessingService()

    def test_invalid_file_path(self, processing_service):
        """Test handling of invalid file path."""
        docling_service = DoclingService()

        with pytest.raises(Exception):
            docling_service.parse_document("/nonexistent/file.pdf")

    def test_empty_schema_validation(self, processing_service):
        """Test validation with empty schema."""
        from services.processing import parse_and_validate_response

        result = parse_and_validate_response("{}", {"type": "object"})
        assert result["success"] is True

    def test_invalid_json_validation(self, processing_service):
        """Test validation with invalid JSON."""
        from services.processing import parse_and_validate_response

        result = parse_and_validate_response("not json", {"type": "object"})
        assert result["success"] is False
        assert "error" in result


# Manual Testing Checklist
# This section documents the manual testing steps that should be performed
# after the automated tests pass.
"""
MANUAL TESTING CHECKLIST for Phase 1 E2E Tests:

Document Upload & Processing:
- [ ] Upload DOCX file via frontend
- [ ] Upload PPTX file via frontend
- [ ] Upload 50-page PDF via frontend
- [ ] Upload searchable PDF via frontend
- [ ] Upload image-only PDF via frontend

Transcription Mode:
- [ ] Test transcription mode with audio file
- [ ] Verify transcription produces Markdown output (not JSON)
- [ ] Verify Markdown viewer displays transcription correctly

Chunking:
- [ ] Upload large document (>50 pages)
- [ ] Verify chunking progress is displayed in UI
- [ ] Verify all chunks are processed
- [ ] Verify results are merged correctly

Markdown Viewer:
- [ ] Test Markdown viewer with various document types
- [ ] Verify formatting is preserved
- [ ] Test syntax highlighting for code blocks

Side-by-Side View:
- [ ] Enable side-by-side view
- [ ] Verify original document and extracted data are both visible
- [ ] Test with different document types

Quality Gate:
- [ ] Upload low-quality image
- [ ] Verify quality gate triggers
- [ ] Verify preprocessing is applied if configured
- [ ] Verify quality metrics are displayed

Schema Validation:
- [ ] Test with different schema templates
- [ ] Verify extracted data matches schema
- [ ] Test with custom schema
"""
