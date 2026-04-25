#!/usr/bin/env python3
"""
Generate test fixtures for OCR Platform integration tests.

This script creates sample documents for testing all extraction methods:
- sample.docx, sample.pptx (docling-parse)
- searchable.pdf, image_only.pdf, large_pdf.pdf (docling-parse / vision / text)
- invoice.pdf (structured extraction target)
- receipt.jpg (image-based vision extraction)
- sample_audio.wav (transcription)
"""

import importlib.util
import struct
import sys
import wave
from pathlib import Path


def check_dependencies():
    missing = []
    if importlib.util.find_spec("docx") is None:
        missing.append("python-docx")
    if importlib.util.find_spec("pptx") is None:
        missing.append("python-pptx")
    if importlib.util.find_spec("reportlab") is None:
        missing.append("reportlab")
    if importlib.util.find_spec("PIL") is None:
        missing.append("Pillow")
    if missing:
        print("Missing required dependencies:")
        for dep in missing:
            print(f"  - {dep}")
        print(f"\nInstall with: pip install {' '.join(missing)}")
        return False
    return True


def create_sample_docx(output_path):
    from docx import Document

    doc = Document()
    doc.add_heading("Quarterly Report Q3 2024", level=1)
    doc.add_paragraph(
        "This report summarizes Q3 2024 performance for Acme Corporation."
    )

    doc.add_heading("Financial Summary", level=2)
    doc.add_paragraph("Total Revenue: $1,250,000")
    doc.add_paragraph("Operating Expenses: $875,000")
    doc.add_paragraph("Net Profit: $375,000")

    doc.add_heading("Key Metrics", level=2)
    doc.add_paragraph("Customer Acquisition: 1,200 new customers")
    doc.add_paragraph("Churn Rate: 2.3%")
    doc.add_paragraph("Average Revenue Per User: $104.17")

    doc.add_heading("Regional Breakdown", level=2)
    doc.add_paragraph("North America: $625,000 (50%)")
    doc.add_paragraph("Europe: $375,000 (30%)")
    doc.add_paragraph("Asia Pacific: $250,000 (20%)")

    doc.save(output_path)
    print(f"Created: {output_path}")


def create_sample_pptx(output_path):
    from pptx import Presentation

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Project Update"
    title_slide.placeholders[1].text = "Q3 2024 Review"

    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    content_slide.shapes.title.text = "Budget Overview"
    body = content_slide.placeholders[1]
    body.text_frame.text = (
        "Allocated: $500,000\nSpent: $425,000\nRemaining: $75,000"
    )

    content_slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    content_slide2.shapes.title.text = "Timeline"
    body2 = content_slide2.placeholders[1]
    body2.text_frame.text = (
        "Phase 1: Complete (Jan-Mar)\nPhase 2: In Progress (Apr-Jun)\nPhase 3: Planned (Jul-Sep)"
    )

    prs.save(output_path)
    print(f"Created: {output_path}")


def create_searchable_pdf(output_path):
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph

    doc = SimpleDocTemplate(str(output_path), pagesize=letter)
    styles = getSampleStyleSheet()
    story = [
        Paragraph("Searchable PDF Document", styles["Heading1"]),
        Paragraph("This is a test paragraph for searchable PDF parsing.", styles["Normal"]),
        Paragraph("Another paragraph with searchable text content.", styles["Normal"]),
        Paragraph("This PDF should be parsed without OCR.", styles["Normal"]),
    ]
    doc.build(story)
    print(f"Created: {output_path}")


def create_image_only_pdf(output_path):
    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("RGB", (800, 600), color="white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 40)
    except Exception:
        try:
            font = ImageFont.truetype("arial.ttf", 40)
        except Exception:
            font = ImageFont.load_default()

    draw.text((50, 50), "Image-Only PDF", fill="black", font=font)
    draw.text((50, 100), "This text is embedded in an image", fill="black", font=font)
    draw.text((50, 150), "OCR is required to extract this text", fill="black", font=font)
    img.save(output_path, "PDF")
    print(f"Created: {output_path}")


def create_large_pdf(output_path, num_pages=50):
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, PageBreak

    doc = SimpleDocTemplate(str(output_path), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    for i in range(num_pages):
        story.append(Paragraph(f"Chapter {i + 1}", styles["Heading1"]))
        story.append(Paragraph(f"This is the content for chapter {i + 1}. " * 20, styles["Normal"]))
        story.append(Paragraph(f"Additional details for section {i + 1}. " * 10, styles["Normal"]))
        if i < num_pages - 1:
            story.append(PageBreak())
    doc.build(story)
    print(f"Created: {output_path} ({num_pages} pages)")


def create_invoice_pdf(output_path):
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors

    doc = SimpleDocTemplate(str(output_path), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph("INVOICE", styles["Heading1"]))
    story.append(Spacer(1, 0.3 * inch))

    info_data = [
        ["Invoice Number:", "INV-2024-0042"],
        ["Date:", "November 15, 2024"],
        ["Due Date:", "December 15, 2024"],
    ]
    info_table = Table(info_data, colWidths=[1.5 * inch, 3 * inch])
    info_table.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    story.append(info_table)
    story.append(Spacer(1, 0.3 * inch))

    story.append(Paragraph("Bill To:", styles["Heading3"]))
    story.append(Paragraph("Acme Corporation\n123 Business Ave\nSan Francisco, CA 94102", styles["Normal"]))
    story.append(Spacer(1, 0.3 * inch))

    story.append(Paragraph("Vendor:", styles["Heading3"]))
    story.append(Paragraph("CloudSync Solutions\n456 Tech Blvd\nAustin, TX 78701", styles["Normal"]))
    story.append(Spacer(1, 0.3 * inch))

    line_data = [
        ["Description", "Qty", "Unit Price", "Total"],
        ["Enterprise License - Annual", "1", "$5,000.00", "$5,000.00"],
        ["Premium Support Package", "1", "$1,200.00", "$1,200.00"],
        ["Data Migration Service", "3", "$400.00", "$1,200.00"],
        ["Training Sessions (8hrs)", "2", "$350.00", "$700.00"],
    ]
    line_table = Table(line_data, colWidths=[2.5 * inch, 0.7 * inch, 1.2 * inch, 1.2 * inch])
    line_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.Color(0.2, 0.4, 0.8)),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ("ALIGN", (1, 0), (-1, -1), "RIGHT"),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    story.append(line_table)
    story.append(Spacer(1, 0.2 * inch))

    totals_data = [
        ["Subtotal:", "$8,100.00"],
        ["Tax (8.25%):", "$668.25"],
        ["Total:", "$8,768.25"],
    ]
    totals_table = Table(totals_data, colWidths=[4.2 * inch, 1.4 * inch])
    totals_table.setStyle(TableStyle([
        ("FONTNAME", (0, -1), (-1, -1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 11),
        ("ALIGN", (1, 0), (1, -1), "RIGHT"),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ("LINEABOVE", (0, -1), (-1, -1), 1, colors.black),
    ]))
    story.append(totals_table)

    doc.build(story)
    print(f"Created: {output_path}")


def create_receipt_image(output_path):
    from PIL import Image, ImageDraw, ImageFont

    width, height = 400, 600
    img = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(img)

    try:
        font_title = ImageFont.truetype("/System/Library/Fonts/Menlo.ttc", 18)
        font_body = ImageFont.truetype("/System/Library/Fonts/Menlo.ttc", 13)
    except Exception:
        try:
            font_title = ImageFont.truetype("courier.ttf", 18)
            font_body = ImageFont.truetype("courier.ttf", 13)
        except Exception:
            font_title = ImageFont.load_default()
            font_body = ImageFont.load_default()

    y = 20
    lines = [
        (font_title, "MARTHA'S CAFE", "center"),
        (font_body, "123 Main Street", "center"),
        (font_body, "Anytown, USA 12345", "center"),
        (font_body, "Tel: (555) 123-4567", "center"),
        (font_body, "", "left"),
        (font_body, "Date: 2024-11-20 14:32", "left"),
        (font_body, "Receipt #: REC-8847", "left"),
        (font_body, "Cashier: Maria S.", "left"),
        (font_body, "", "left"),
        (font_body, "-" * 36, "left"),
        (font_body, "QTY  ITEM              AMOUNT", "left"),
        (font_body, "-" * 36, "left"),
        (font_body, " 2   Espresso           $8.00", "left"),
        (font_body, " 1   Croissant          $4.50", "left"),
        (font_body, " 1   Avocado Toast     $12.00", "left"),
        (font_body, " 1   Fresh Juice        $6.50", "left"),
        (font_body, "-" * 36, "left"),
        (font_body, "    Subtotal:          $31.00", "left"),
        (font_body, "    Tax (9.5%):         $2.95", "left"),
        (font_body, "    Tip:                $5.00", "left"),
        (font_body, " " * 36, "left"),
        (font_title, "TOTAL:  $38.95", "left"),
        (font_body, "", "left"),
        (font_body, "Payment: Visa ****4829", "left"),
        (font_body, "", "left"),
        (font_body, "Thank you for visiting!", "center"),
    ]
    for font, text, align in lines:
        if align == "center":
            bbox = draw.textbbox((0, 0), text, font=font)
            tw = bbox[2] - bbox[0]
            x = (width - tw) // 2
        else:
            x = 20
        draw.text((x, y), text, fill="black", font=font)
        y += 20 if font == font_body else 26

    img.save(output_path, "JPEG", quality=95)
    print(f"Created: {output_path}")


def create_sample_audio(output_path):
    sample_rate = 22050
    duration = 2
    num_samples = sample_rate * duration

    with wave.open(str(output_path), "w") as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(sample_rate)
        frames = struct.pack("<" + "h" * num_samples, *([0] * num_samples))
        wf.writeframes(frames)
    print(f"Created: {output_path} ({duration}s silence)")


def create_multi_page_searchable_pdf(output_path):
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, PageBreak

    doc = SimpleDocTemplate(str(output_path), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    pages = [
        ("Contract Agreement", "This Master Service Agreement is entered into on January 10, 2025 between TechCorp Inc. (Client) and DataFlow Solutions (Provider)."),
        ("Scope of Work", "The Provider shall deliver data pipeline automation services including ETL development, API integration, and real-time monitoring dashboards."),
        ("Payment Terms", "Total contract value: $45,000. Payment schedule: 30% upon signing, 40% at midpoint review, 30% upon completion. Net 30 payment terms."),
    ]
    for title, content in pages:
        story.append(Paragraph(title, styles["Heading1"]))
        story.append(Paragraph(content, styles["Normal"]))
        story.append(PageBreak())
    story.pop()

    doc.build(story)
    print(f"Created: {output_path} ({len(pages)} pages)")


def main():
    print("Generating OCR Platform integration test fixtures...")
    if not check_dependencies():
        sys.exit(1)

    fixtures_dir = Path(__file__).parent
    fixtures_dir.mkdir(parents=True, exist_ok=True)

    generators = [
        ("sample.docx", lambda p: create_sample_docx(p)),
        ("sample.pptx", lambda p: create_sample_pptx(p)),
        ("searchable.pdf", lambda p: create_searchable_pdf(p)),
        ("image_only.pdf", lambda p: create_image_only_pdf(p)),
        ("large_pdf.pdf", lambda p: create_large_pdf(p)),
        ("invoice.pdf", lambda p: create_invoice_pdf(p)),
        ("receipt.jpg", lambda p: create_receipt_image(p)),
        ("sample_audio.wav", lambda p: create_sample_audio(p)),
        ("multi_page.pdf", lambda p: create_multi_page_searchable_pdf(p)),
    ]

    try:
        for name, gen in generators:
            gen(fixtures_dir / name)

        print("\nAll fixtures generated successfully!")
        print(f"Fixture directory: {fixtures_dir}")
        print("\nGenerated files:")
        for fp in sorted(fixtures_dir.iterdir()):
            if fp.is_file() and fp.suffix in (".pdf", ".docx", ".pptx", ".jpg", ".wav", ".py"):
                print(f"  {fp.name:30s} {fp.stat().st_size:>10,} bytes")
    except Exception as e:
        print(f"\nError generating fixtures: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()

